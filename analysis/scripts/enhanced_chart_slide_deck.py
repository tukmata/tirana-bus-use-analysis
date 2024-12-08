import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO

# Load data
data_file = "aggregated_data.csv"
data = pd.read_csv(data_file)

# Create a presentation
prs = Presentation()

# Helper function to add a slide with title and content
def add_title_content_slide(prs, title_text, content_text):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = title_text
    content.text = content_text

# Helper function to add a table slide
def add_table_slide(prs, title_text, dataframe):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = title_text

    rows, cols = dataframe.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(1), Inches(1.5), Inches(8), Inches(3)).table

    # Add column headers
    for col_idx, col_name in enumerate(dataframe.columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.text_frame.paragraphs[0].font.bold = True

    # Add data
    for row_idx, row in dataframe.iterrows():
        for col_idx, value in enumerate(row):
            table.cell(row_idx + 1, col_idx).text = str(value)

# Helper function to add a chart slide
def add_chart_slide(prs, title_text, chart_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = title_text
    left = Inches(1)
    top = Inches(1.5)
    slide.shapes.add_picture(chart_path, left, top, width=Inches(8), height=Inches(4))

# Function to create a chart
def create_chart(data, x_col, y_col, title, chart_path):
    plt.figure(figsize=(10, 6))
    plt.bar(data[x_col], data[y_col], color="skyblue")
    plt.title(title, fontsize=14)
    plt.xlabel(x_col, fontsize=12)
    plt.ylabel(y_col, fontsize=12)
    plt.xticks(rotation=45, fontsize=10)
    plt.tight_layout()
    plt.savefig(chart_path)
    plt.close()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Tirana Bus Usage Analysis"
subtitle.text = "Exploring Usage Patterns, Demographics, and Time Slots for Better Urban Planning"

# Slide 2: Overview and Objectives
add_title_content_slide(
    prs,
    "Overview and Objectives",
    (
        "This presentation explores bus usage patterns in Tirana, focusing on:\n"
        "- Revenue generation by age groups, time slots, and bus lines\n"
        "- Identifying peak times and high-performing bus lines\n"
        "- Recommendations to optimize services and enhance commuter experience"
    ),
)

# Slide 3: Revenue Breakdown by Age Group
age_group_revenue = data.groupby("age_group")["total_revenue"].sum().reset_index()
chart_path = "age_group_revenue.png"
create_chart(age_group_revenue, "age_group", "total_revenue", "Revenue by Age Group", chart_path)
add_chart_slide(prs, "Revenue Breakdown by Age Group", chart_path)

# Slide 4: Revenue Breakdown by Time Slot
time_slot_revenue = data.groupby("time_slot")["total_revenue"].sum().reset_index()
chart_path = "time_slot_revenue.png"
create_chart(time_slot_revenue, "time_slot", "total_revenue", "Revenue by Time Slot", chart_path)
add_chart_slide(prs, "Revenue Breakdown by Time Slot", chart_path)

# Slide 5: Top Performing Bus Lines
bus_line_revenue = data.groupby("bus_line")["total_revenue"].sum().reset_index().sort_values(by="total_revenue", ascending=False).head(5)
chart_path = "top_bus_lines.png"
create_chart(bus_line_revenue, "bus_line", "total_revenue", "Top Performing Bus Lines", chart_path)
add_chart_slide(prs, "Top Performing Bus Lines", chart_path)

# Slide 6: Recommendations
add_title_content_slide(
    prs,
    "Recommendations",
    (
        "Based on the analysis, we recommend:\n"
        "- Targeted marketing campaigns for age group 19-60\n"
        "- Increased bus frequency during peak hours (07:00-10:00)\n"
        "- Enhanced capacity and services for Bus Line 22\n"
        "- Implementing a dynamic pricing model to encourage off-peak travel"
    ),
)

# Slide 7: Conclusion
add_title_content_slide(
    prs,
    "Conclusion",
    (
        "This analysis provides actionable insights for improving Tirana's public transportation system:\n"
        "- Enhanced commuter experience through optimized services\n"
        "- Improved revenue generation and operational efficiency\n"
        "- Data-driven planning for sustainable urban mobility"
    ),
)

# Save presentation
output_file = "Tirana_Bus_Usage_Analysis_Enhanced.pptx"
prs.save(output_file)
print(f"Enhanced presentation saved as {output_file}.")
