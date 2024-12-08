import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load data
data_file = "aggregated_data.csv"
data = pd.read_csv(data_file)

# Create a presentation
prs = Presentation()

# Function to add a title and content slide
def add_title_content_slide(prs, title_text, content_text):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = title_text
    content.text = content_text

# Function to add a table slide
def add_table_slide(prs, title_text, dataframe):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = title_text

    rows, cols = dataframe.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(1), Inches(1.5), Inches(8), Inches(3)).table

    # Add column headers
    for col_idx, col_name in enumerate(dataframe.columns):
        cell = table.cell(0, col_idx)
        cell.text = col_name
        cell.text_frame.paragraphs[0].font.bold = True

    # Add data
    for row_idx, row in dataframe.iterrows():
        for col_idx, value in enumerate(row):
            table.cell(row_idx + 1, col_idx).text = str(value)

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Tirana Bus Usage Analysis"
subtitle.text = "Exploring Demographics, Time Slots, and Bus Lines"

# Slide 2: Key Insights
add_title_content_slide(
    prs,
    "Key Insights",
    (
        "• Highest revenue generated by passengers aged 19-60.\n"
        "• Peak time slot: 07:00-10:00.\n"
        "• Bus Line 22 contributes most to total revenue.\n"
    ),
)

# Slide 3: Revenue Breakdown by Age Group
age_group_revenue = data.groupby("age_group")["total_revenue"].sum().reset_index()
add_table_slide(prs, "Revenue Breakdown by Age Group", age_group_revenue)

# Slide 4: Revenue Breakdown by Time Slot
time_slot_revenue = data.groupby("time_slot")["total_revenue"].sum().reset_index()
add_table_slide(prs, "Revenue Breakdown by Time Slot", time_slot_revenue)

# Slide 5: Recommendations
add_title_content_slide(
    prs,
    "Recommendations",
    (
        "• Target marketing campaigns for age group 19-60.\n"
        "• Optimize services during peak times (07:00-10:00).\n"
        "• Focus on improving capacity for Bus Line 22.\n"
    ),
)

# Save presentation
output_file = "Enhanced_Tirana_Bus_Usage_Analysis.pptx"
prs.save(output_file)
print(f"Presentation saved as {output_file}.")